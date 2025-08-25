"""
PPTX Image Extraction and OCR Processing Module

This module handles PowerPoint file processing, image extraction, and OCR for the IDC RAG system.
Specifically designed to extract partner logos and company information from presentation slides.
"""

import os
import io
import json
from typing import List, Dict, Optional, Tuple
from pathlib import Path

# Core libraries
from pptx import Presentation
from PIL import Image, ImageEnhance, ImageFilter
import cv2
import numpy as np

# OCR and configuration
import pytesseract
from config import PPTX_PROCESSING, OCR_SETTINGS

# Import logging configuration
from logging_config import (
    setup_clean_logging, pptx_logger, log_processing_summary, log_skipped_formats
)

# Initialize clean logging
setup_clean_logging()


class PPTXImageExtractor:
    """
    Extracts images from PowerPoint presentations and prepares them for OCR processing.
    """
    
    def __init__(self):
        self.config = PPTX_PROCESSING
        self.ocr_config = OCR_SETTINGS
        print("PPTX Image Extractor initialized")
    
    def extract_images_from_pptx(self, pptx_path: str) -> Tuple[List[Dict], Dict]:
        """
        Extract all images from a PowerPoint presentation.
        
        Args:
            pptx_path: Path to the PPTX file
            
        Returns:
            Tuple of (extracted images list, processing stats)
        """
        try:
            presentation = Presentation(pptx_path)
            extracted_images = []
            stats = {'total_found': 0, 'skipped_formats': [], 'extracted': 0}
            
            for slide_num, slide in enumerate(presentation.slides, 1):
                slide_images, slide_stats = self._extract_slide_images(slide, slide_num, pptx_path)
                extracted_images.extend(slide_images)
                stats['total_found'] += slide_stats['total_found']
                stats['skipped_formats'].extend(slide_stats['skipped_formats'])
                stats['extracted'] += len(slide_images)
            
            # Log clean summary instead of individual messages
            if stats['skipped_formats']:
                log_skipped_formats(os.path.basename(pptx_path), 
                                  len(stats['skipped_formats']), 
                                  stats['skipped_formats'])
            
            return extracted_images, stats
            
        except Exception as e:
            pptx_logger.error(f"Error extracting images from {pptx_path}: {e}")
            return [], {'total_found': 0, 'skipped_formats': [], 'extracted': 0}
    
    def _extract_slide_images(self, slide, slide_num: int, source_file: str) -> Tuple[List[Dict], Dict]:
        """
        Extract images from a specific slide.
        
        Args:
            slide: PowerPoint slide object
            slide_num: Slide number
            source_file: Source PPTX file path
            
        Returns:
            Tuple of (image data list, processing stats)
        """
        slide_images = []
        stats = {'total_found': 0, 'skipped_formats': []}
        
        for shape_idx, shape in enumerate(slide.shapes):
            if hasattr(shape, 'image') and shape.image:
                stats['total_found'] += 1
                
                try:
                    # Get image data
                    image_data = shape.image.blob
                    image_format = shape.image.ext.lower()
                    
                    # Skip unsupported formats (WMF, EMF) silently
                    unsupported_formats = ['.wmf', '.emf']
                    if image_format in unsupported_formats:
                        stats['skipped_formats'].append(image_format)
                        continue
                    
                    # Only process supported formats
                    supported_formats = ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff']
                    if image_format not in supported_formats:
                        stats['skipped_formats'].append(image_format)
                        continue
                    
                    # Create PIL Image
                    pil_image = Image.open(io.BytesIO(image_data))
                    
                    # Convert palette images with transparency to RGBA to avoid warnings
                    if pil_image.mode == 'P' and 'transparency' in pil_image.info:
                        pil_image = pil_image.convert('RGBA')
                    
                    # Resize if too large
                    if self.config['max_image_size']:
                        pil_image = self._resize_image(pil_image, self.config['max_image_size'])
                    
                    # Create metadata
                    metadata = {
                        'source_file': source_file,
                        'slide_number': slide_num,
                        'shape_index': shape_idx,
                        'image_format': image_format,
                        'image_size': pil_image.size,
                        'extraction_method': 'pptx_shape'
                    }
                    
                    slide_images.append({
                        'image': pil_image,
                        'metadata': metadata,
                        'raw_data': image_data
                    })
                    
                except Exception as e:
                    # Only log actual errors, not format issues
                    if "cannot find loader" not in str(e).lower():
                        pptx_logger.warning(f"Could not process image from slide {slide_num}, shape {shape_idx}")
        
        return slide_images, stats
    
    def _resize_image(self, image: Image.Image, max_size: Tuple[int, int]) -> Image.Image:
        """
        Resize image while maintaining aspect ratio.
        
        Args:
            image: PIL Image object
            max_size: Maximum (width, height)
            
        Returns:
            Resized PIL Image
        """
        image.thumbnail(max_size, Image.Resampling.LANCZOS)
        return image
    
    def preprocess_image_for_ocr(self, pil_image: Image.Image) -> Image.Image:
        """
        Preprocess image to improve OCR accuracy.
        
        Args:
            pil_image: Input PIL Image
            
        Returns:
            Preprocessed PIL Image
        """
        if not self.ocr_config['preprocess_images']:
            return pil_image
        
        try:
            # Convert to RGB if needed
            if pil_image.mode != 'RGB':
                pil_image = pil_image.convert('RGB')
            
            # Enhance contrast
            if self.ocr_config['enhance_contrast']:
                enhancer = ImageEnhance.Contrast(pil_image)
                pil_image = enhancer.enhance(1.5)
            
            # Convert to OpenCV format for advanced preprocessing
            opencv_image = cv2.cvtColor(np.array(pil_image), cv2.COLOR_RGB2BGR)
            
            # Convert to grayscale
            gray = cv2.cvtColor(opencv_image, cv2.COLOR_BGR2GRAY)
            
            # Remove noise
            if self.ocr_config['remove_noise']:
                gray = cv2.medianBlur(gray, 3)
            
            # Apply adaptive threshold
            processed = cv2.adaptiveThreshold(
                gray, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY, 11, 2
            )
            
            # Convert back to PIL
            return Image.fromarray(processed)
            
        except Exception as e:
            # Suppress common format errors
            if "cannot find loader" not in str(e).lower():
                print(f"Warning: Image preprocessing failed, using original image")
            return pil_image


class PPTXOCRProcessor:
    """
    Performs OCR on extracted images using pytesseract.
    """
    
    def __init__(self):
        self.config = OCR_SETTINGS
        print("PPTX OCR Processor initialized")
    
    def extract_text_from_image(self, pil_image: Image.Image, metadata: Dict) -> Dict:
        """
        Extract text from an image using OCR.
        
        Args:
            pil_image: PIL Image object
            metadata: Image metadata
            
        Returns:
            Dictionary with extracted text and confidence
        """
        try:
            # Get OCR configuration
            custom_config = self.config['tesseract_config']
            
            # Extract text with confidence data
            ocr_data = pytesseract.image_to_data(
                pil_image, 
                config=custom_config, 
                output_type=pytesseract.Output.DICT
            )
            
            # Filter by confidence
            extracted_text = []
            confidences = []
            
            for i in range(len(ocr_data['text'])):
                text = ocr_data['text'][i].strip()
                conf = int(ocr_data['conf'][i])
                
                if text and conf >= self.config['min_confidence']:
                    extracted_text.append(text)
                    confidences.append(conf)
            
            # Combine text
            full_text = ' '.join(extracted_text)
            avg_confidence = np.mean(confidences) if confidences else 0
            
            return {
                'extracted_text': full_text,
                'confidence': avg_confidence,
                'word_count': len(extracted_text),
                'method': 'pytesseract',
                'metadata': metadata
            }
            
        except Exception as e:
            print(f"Error during OCR processing: {e}")
            return {
                'extracted_text': '',
                'confidence': 0,
                'word_count': 0,
                'method': 'pytesseract',
                'error': str(e),
                'metadata': metadata
            }


class PPTXProcessor:
    """
    Main PPTX processing class that combines image extraction and OCR.
    """
    
    def __init__(self):
        self.image_extractor = PPTXImageExtractor()
        self.ocr_processor = PPTXOCRProcessor()
        print("PPTX Processor initialized successfully")
    
    def process_pptx_file(self, pptx_path: str) -> Dict:
        """
        Process a complete PPTX file and extract all text content.
        
        Args:
            pptx_path: Path to the PPTX file
            
        Returns:
            Dictionary with extracted content and metadata
        """
        try:
            # Extract presentation metadata
            presentation = Presentation(pptx_path)
            slide_count = len(presentation.slides)
            
            # Extract images
            extracted_images, extraction_stats = self.image_extractor.extract_images_from_pptx(pptx_path)
            
            # Process each image with OCR
            ocr_results = []
            for image_data in extracted_images:
                # Preprocess image
                preprocessed_image = self.image_extractor.preprocess_image_for_ocr(
                    image_data['image']
                )
                
                # Extract text
                ocr_result = self.ocr_processor.extract_text_from_image(
                    preprocessed_image, image_data['metadata']
                )
                
                if ocr_result['extracted_text']:  # Only include non-empty results
                    ocr_results.append(ocr_result)
            
            # Log clean processing summary
            log_processing_summary(
                os.path.basename(pptx_path),
                extraction_stats['total_found'],
                len(ocr_results),
                "PPTX+OCR"
            )
            
            # Compile final result
            result = {
                'source_file': pptx_path,
                'slide_count': slide_count,
                'images_processed': len(extracted_images),
                'successful_ocr': len(ocr_results),
                'ocr_results': ocr_results,
                'processing_method': 'pptx_ocr'
            }
            
            return result
            
        except Exception as e:
            print(f"Error processing PPTX file {pptx_path}: {e}")
            return {
                'source_file': pptx_path,
                'error': str(e),
                'processing_method': 'pptx_ocr'
            }


# Utility functions for integration
def extract_text_from_pptx(pptx_path: str) -> str:
    """
    Extract all text content from a PPTX file for integration with existing data loader.
    
    Args:
        pptx_path: Path to PPTX file
        
    Returns:
        Combined extracted text
    """
    processor = PPTXProcessor()
    result = processor.process_pptx_file(pptx_path)
    
    if 'error' in result:
        return ""
    
    # Combine all OCR results
    all_text = []
    for ocr_result in result.get('ocr_results', []):
        text = ocr_result['extracted_text']
        metadata = ocr_result['metadata']
        slide_num = metadata['slide_number']
        
        # Format with slide context
        formatted_text = f"--- Slide {slide_num} (OCR) ---\n{text}\n\n"
        all_text.append(formatted_text)
    
    return ''.join(all_text).strip()


if __name__ == "__main__":
    # Test the processor
    processor = PPTXProcessor()
    
    # Example usage
    test_file = "/path/to/test.pptx"
    if os.path.exists(test_file):
        result = processor.process_pptx_file(test_file)
        print(json.dumps(result, indent=2, default=str))
    else:
        print("No test file found. PPTX processor is ready for integration.")